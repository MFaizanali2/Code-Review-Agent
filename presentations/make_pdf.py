from pptx import Presentation
from fpdf import FPDF
import os, re

prs = Presentation(os.path.join(os.path.dirname(__file__), "Code-Review-Agent-Presentation.pptx"))

def clean(s):
    s = s.replace("\u2014", "-").replace("\u2013", "-")
    s = s.replace("\u2018", "'").replace("\u2019", "'")
    s = s.replace("\u201c", '"').replace("\u201d", '"')
    s = s.replace("\u2022", "-").replace("\u2026", "...")
    s = re.sub(r"[^\x00-\x7F]", "", s)
    return s.strip()

pdf = FPDF(orientation="L", format="A4")
pdf.set_auto_page_break(auto=True, margin=20)

BLUE = (13, 27, 42)
ACCENT = (46, 134, 193)
GOLD = (243, 156, 18)
WHITE = (255, 255, 255)
LGRAY = (200, 210, 220)
FONT = "Helvetica"

for i, slide in enumerate(prs.slides, 1):
    pdf.add_page()
    pdf.set_fill_color(*BLUE)
    pdf.rect(0, 0, 297, 210, "F")
    pdf.set_fill_color(*ACCENT)
    pdf.rect(0, 0, 3, 210, "F")

    shapes = slide.shapes
    texts = [clean(s.text) for s in shapes if s.has_text_frame and s.text.strip()]
    title = texts[0] if texts else "(no title)"
    body = texts[1:] if len(texts) > 1 else []
    notes = ""
    if slide.has_notes_slide:
        notes = clean(slide.notes_slide.notes_text_frame.text)

    pdf.set_text_color(*WHITE)
    pdf.set_font(FONT, "B", 20)
    pdf.set_xy(10, 10)
    pdf.cell(0, 12, f"Slide {i}: {title[:80]}")
    pdf.set_draw_color(*ACCENT)
    pdf.set_line_width(0.5)
    pdf.line(10, 25, 287, 25)

    pdf.set_text_color(*LGRAY)
    pdf.set_font(FONT, "", 12)
    y = 30
    for line in body:
        if line:
            pdf.set_xy(12, y)
            pdf.cell(0, 7, line[:130])
            y += 8

    if notes:
        if y < 145:
            y = 145
        pdf.set_fill_color(22, 42, 64)
        pdf.rect(10, y, 277, min(55, 190 - y), "F")
        pdf.set_text_color(*GOLD)
        pdf.set_font(FONT, "B", 10)
        pdf.set_xy(14, y + 3)
        pdf.cell(0, 5, "SPEAKER NOTES:")
        pdf.set_text_color(*LGRAY)
        pdf.set_font(FONT, "", 9)
        pdf.set_xy(14, y + 10)
        pdf.multi_cell(268, 4.5, notes[:500])

out = os.path.join(os.path.dirname(__file__), "Code-Review-Agent-Presentation.pdf")
pdf.output(out)
print(f"PDF saved: {out} ({os.path.getsize(out)} bytes, {len(prs.slides)} slides)")
