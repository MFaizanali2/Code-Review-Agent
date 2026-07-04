from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE_DARK = RGBColor(0x0D, 0x1B, 0x2A)
BLUE_MID = RGBColor(0x1B, 0x3A, 0x5C)
BLUE_ACCENT = RGBColor(0x2E, 0x86, 0xC1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE6, 0xED)
GOLD = RGBColor(0xF3, 0x9C, 0x12)
GREEN = RGBColor(0x27, 0xAE, 0x60)
DARK_TEXT = RGBColor(0x1C, 0x2D, 0x3D)

def add_bg(slide, color=BLUE_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, title_text, bullets, notes_text="", accent_color=BLUE_ACCENT):
    add_bg(slide, BLUE_DARK)
    acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), accent_color)
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), title_text, 36, WHITE, True)
    line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), accent_color)
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(20)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(12)
        p.level = 0
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

def add_card(slide, left, top, width, height, title, body, title_color=BLUE_ACCENT, body_color=LIGHT_GRAY):
    card = add_rect(slide, left, top, width, height, BLUE_MID, RGBColor(0x2C, 0x4A, 0x6F))
    card.shadow.inherit = False
    add_textbox(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5), title, 22, title_color, True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), height - Inches(0.8), body, 16, body_color)

def add_arrow_right(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.6), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape

def add_cycle_box(slide, left, top, width, height, label, desc, color=BLUE_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Calibri"
    tf.paragraphs[0].space_before = Pt(12)
    return shape

def add_image_placeholder(slide, left, top, width, height, label="[Screenshot]"):
    shape = add_rect(slide, left, top, width, height, RGBColor(0x24, 0x3B, 0x55), RGBColor(0x3A, 0x5A, 0x7A))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x8A, 0xAA, 0xC8)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(30)
    return shape

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2), "Code Review Agent", 54, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8), "AI-Powered Automated Code Analysis", 28, BLUE_ACCENT, False, PP_ALIGN.CENTER)
line = add_rect(slide, Inches(4.5), Inches(3.6), Inches(4.3), Inches(0.04), GOLD)
add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6), "An Agentic AI Learning Project — ReACT Loop Architecture", 20, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5), "Team: Person 1 · Person 2 · Person 3 · Person 4 · Person 5", 18, RGBColor(0x8A, 0xAA, 0xC8), False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5), "January 2025", 16, RGBColor(0x6A, 0x8A, 0xA8), False, PP_ALIGN.CENTER)
notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (1 min): Welcome everyone! Today we'll present our Code Review Agent — an AI-powered system that automatically analyzes code quality. Built over 2 weeks by a team of 5, this project demonstrates agentic AI using the ReACT loop pattern. We achieved an 85/100 audit score with 257 passing tests. Let's dive in!"

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: PROJECT OVERVIEW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Project Overview", [
    "What: AI system that autonomously reviews code from GitHub repositories",
    "Why: Learn Agentic AI development with ReACT loop architecture",
    "How: Agent reasons (Think → Act → Observe → Reflect) using LLM + Tools",
    "Status: Complete • 85/100 audit score • 257 tests passing • Production-ready",
    "Team: 5 members • 2-week sprint • 6,000+ lines of Python code",
], "SPEAKER NOTES (1.5 min): This project was designed to solve real code review bottlenecks... We chose the ReACT (Reasoning + Acting) loop because it's the industry standard for autonomous AI agents. Key result: 85/100 audit score shows high code quality.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "The Problem", [
    "Manual code reviews are time-consuming and inconsistent between reviewers",
    "Teams struggle to catch security vulnerabilities, performance issues, and bugs",
    "Small teams lack bandwidth for thorough reviews of every pull request",
    "Need for an automated, intelligent system that learns and adapts to codebases",
], "SPEAKER NOTES (1 min): Manual code review is essential but expensive... A single review can take 30-60 minutes. Our agent automates the initial pass, catching common issues instantly.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: SOLUTION OVERVIEW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "Solution: Code Review Agent", 36, WHITE, True)
line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(1.5), Inches(6), Inches(0.5), "End-to-End Automated Code Analysis Pipeline", 20, LIGHT_GRAY)

# Flow boxes
boxes = [
    (Inches(0.6), Inches(2.4), "GitHub URL", "User submits a\nrepository URL for review"),
    (Inches(4.0), Inches(2.4), "Agent Core", "ReACT loop processes\ncode with LLM reasoning"),
    (Inches(7.4), Inches(2.4), "Analysis Tools", "5+ specialized tools\nanalyze the codebase"),
    (Inches(10.8), Inches(2.4), "Report", "Comprehensive review\nwith scores & issues"),
]
colors = [GREEN, BLUE_ACCENT, GOLD, GREEN]
for (l, t, title, desc), c in zip(boxes, colors):
    add_cycle_box(slide, l, t, Inches(2.8), Inches(1.8), title, desc, c)

arrows = [(Inches(3.4), Inches(3.1)), (Inches(6.8), Inches(3.1)), (Inches(10.2), Inches(3.1))]
for l, t in arrows:
    add_arrow_right(slide, l, t)

add_textbox(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.4), "Tech Stack: Python 3.10+ • FastAPI • Streamlit • Gemini AI • SQLAlchemy • Pytest", 16, RGBColor(0x8A, 0xAA, 0xC8), False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4), "257 Tests ✅  •  85/100 Audit Score  •  6,000+ Lines of Code  •  2-Week Sprint", 16, GOLD, False, PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (1.5 min): Our solution is a complete pipeline from GitHub URL to comprehensive review report. The user submits a repository URL, our agent fetches the code, analyzes it using 5+ tools, and generates a detailed report with scores, issues found, and improvement suggestions."

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: THE ReACT LOOP
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "What is the ReACT Loop?", 36, WHITE, True)
line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5), "Reasoning + Acting = Intelligent Agent Behavior", 20, LIGHT_GRAY)

# 4 boxes in a cycle layout
cycle_items = [
    (Inches(1.0), Inches(2.3), "🧠  Think", "Agent analyzes context\ndecides next action", RGBColor(0x2E, 0x86, 0xC1)),
    (Inches(4.6), Inches(2.3), "⚡  Act", "Execute selected tool\n(run analysis, fetch data)", RGBColor(0x27, 0xAE, 0x60)),
    (Inches(8.2), Inches(2.3), "👁️  Observe", "Collect tool output\n& environment feedback", RGBColor(0xF3, 0x9C, 0x12)),
    (Inches(11.8), Inches(2.3), "🔄  Reflect", "Process observations\nplan next iteration", RGBColor(0xE7, 0x4C, 0x3C)),
]
for l, t, label, desc, c in cycle_items:
    add_cycle_box(slide, l, t, Inches(2.3), Inches(2.2), label, desc, c)

# Cycle arrows (top row)
cycle_arrows = [(Inches(3.3), Inches(3.0)), (Inches(6.9), Inches(3.0)), (Inches(10.5), Inches(3.0))]
for l, t in cycle_arrows:
    add_arrow_right(slide, l, t)

# Loop back arrow
add_textbox(slide, Inches(4.5), Inches(5.2), Inches(4), Inches(0.4), "↻ Loop continues until task complete or max iterations reached", 16, GOLD, False, PP_ALIGN.CENTER)

# Bottom text
add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.5), "The ReACT loop enables the agent to reason about code, take actions like fetching files or running analysis, observe results, and iteratively refine its approach — exactly like a human reviewer would.", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (2 min): Let me explain ReACT — it stands for Reasoning + Acting. Unlike traditional chatbots that just respond, ReACT agents actively think about what to do, take actions, observe results, and reflect on next steps. This is why modern AI agents work so well..."

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: KEY COMPONENTS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "Key Components", 36, WHITE, True)
line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), BLUE_ACCENT)

comps = [
    ("1. Agent Core", "ReACT loop controller — orchestrates Think → Act → Observe → Reflect cycle\nManages state, context window, iteration limits, and tool routing"),
    ("2. Analysis Tools", "5+ specialized tools: Code fetch, structure analysis, security audit,\nperformance check, dependency analysis, quality scoring"),
    ("3. LLM Integration", "Gemini 2.0 Flash (primary) • OpenAI GPT-4o-mini (alternative)\nHandles reasoning, tool selection, report generation"),
    ("4. FastAPI Backend", "REST API endpoints for review submission, history, health checks\nSQLAlchemy + SQLite for persistence • async support"),
    ("5. Streamlit UI", "Web interface with code input, review history, live progress\nSidebar configuration • real-time status updates"),
]
for i, (title, body) in enumerate(comps):
    top = Inches(1.5 + i * 1.15)
    add_textbox(slide, Inches(0.6), top, Inches(3.5), Inches(0.4), title, 22, BLUE_ACCENT, True)
    add_textbox(slide, Inches(4.2), top, Inches(8.5), Inches(1.0), body, 16, LIGHT_GRAY)

notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (2 min): The system has 5 main components working together... The Agent Core orchestrates everything. Tools give the agent capabilities. LLM provides reasoning. FastAPI serves the API. Streamlit provides the UI."

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: TECHNICAL ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "Technical Architecture", 36, WHITE, True)
line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), BLUE_ACCENT)

# Backend box
add_rect(slide, Inches(0.6), Inches(1.6), Inches(7.5), Inches(5.2), RGBColor(0x14, 0x2B, 0x44), RGBColor(0x2C, 0x4A, 0x6F))
add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.4), "BACKEND (FastAPI Server)", 20, GOLD, True)

# Sub-components in backend
backend_comps = [
    (Inches(0.8), Inches(2.3), "API Layer\nroutes.py, dependencies.py", GREEN),
    (Inches(0.8), Inches(3.5), "Agent Core\nagent.py, agent_types.py", BLUE_ACCENT),
    (Inches(3.0), Inches(2.3), "Analysis Tools\ncode_fetch, security,\nperformance, quality", GOLD),
    (Inches(3.0), Inches(3.5), "LLM Integration\ngemini_client, openai_client,\nllm_client", RGBColor(0xE7, 0x4C, 0x3C)),
    (Inches(5.2), Inches(2.3), "Database\nSQLAlchemy + SQLite\nReviewRecord model", GREEN),
    (Inches(5.2), Inches(3.5), "Config & Utils\nsettings.py, logger.py\nvalidation", RGBColor(0x8E, 0x44, 0xAD)),
]
for l, t, label, c in backend_comps:
    s = add_rect(slide, l, t, Inches(2.0), Inches(1.0), c, RGBColor(0xFF, 0xFF, 0xFF))
    s.fill.solid()
    s.fill.fore_color.rgb = c
    txb = add_textbox(slide, l + Inches(0.1), t + Inches(0.05), Inches(1.8), Inches(0.9), label, 11, WHITE, False)

# Frontend box
add_rect(slide, Inches(8.5), Inches(1.6), Inches(4.3), Inches(2.5), RGBColor(0x14, 0x2B, 0x44), RGBColor(0x2C, 0x4A, 0x6F))
add_textbox(slide, Inches(8.7), Inches(1.7), Inches(4), Inches(0.4), "FRONTEND (Streamlit)", 20, GOLD, True)
frontend_comps = [
    (Inches(8.7), Inches(2.3), "UI Components\nstreamlit_app.py", GREEN),
    (Inches(10.5), Inches(2.3), "API Client\napi_client.py", BLUE_ACCENT),
]
for l, t, label, c in frontend_comps:
    add_rect(slide, l, t, Inches(1.6), Inches(0.9), c)
    txb = add_textbox(slide, l + Inches(0.08), t + Inches(0.05), Inches(1.5), Inches(0.8), label, 11, WHITE, False)

# Arrow between frontend and backend
add_arrow_right(slide, Inches(8.3), Inches(2.6))

# Testing box
add_rect(slide, Inches(8.5), Inches(4.5), Inches(4.3), Inches(2.3), RGBColor(0x14, 0x2B, 0x44), RGBColor(0x2C, 0x4A, 0x6F))
add_textbox(slide, Inches(8.7), Inches(4.6), Inches(4), Inches(0.4), "TESTING (Pytest)", 20, GOLD, True)
add_textbox(slide, Inches(8.7), Inches(5.2), Inches(4), Inches(1.2), "• 257 tests (unit + integration)\n• 85%+ code coverage\n• pytest-asyncio for async\n• Mock LLM provider for CI\n• SQLite in-memory for tests", 14, LIGHT_GRAY)

notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (2 min): Here's the full architecture... Backend has 6 main areas: API, Agent, Tools, LLM, Database, and Config. Frontend is Streamlit with API client. Testing has 257 tests with 85%+ coverage using mock providers."

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: TEAM CONTRIBUTIONS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "Team Contributions", 36, WHITE, True)
line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), BLUE_ACCENT)

members = [
    ("Person 1", "Agent Architect", "Designed ReACT loop architecture\nImplemented agent core, context management\nState machine, iteration control", BLUE_ACCENT),
    ("Person 2", "Tool Engineer", "Built 5+ analysis tools\nCode fetch, security audit, performance\nFile structure analyzer, quality scorer", GREEN),
    ("Person 3", "LLM Specialist", "Gemini 2.0 Flash integration\nOpenAI provider implementation\nPrompt engineering, response parsing", GOLD),
    ("Person 4", "Backend Developer", "FastAPI routes, dependencies\nDatabase models, SQLAlchemy\nConfiguration, logging utilities", RGBColor(0xE7, 0x4C, 0x3C)),
    ("Person 5", "Frontend Developer", "Streamlit UI with 4 screens\nAPI client integration\nProgress tracking, history view", RGBColor(0x8E, 0x44, 0xAD)),
]
for i, (name, role, desc, color) in enumerate(members):
    left = Inches(0.6 + i * 2.5)
    add_card(slide, left, Inches(1.6), Inches(2.3), Inches(5.2), f"{name}\n{role}", desc, color)

notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (2 min): Each team member owned a critical component... Person 1 built the ReACT loop, Person 2 created analysis tools, Person 3 handled LLM integration, Person 4 built the backend API, Person 5 created the Streamlit frontend."

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: FEATURES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "Features Implemented", 36, WHITE, True)
line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), BLUE_ACCENT)

features = [
    "✅ ReACT Loop — Autonomous reasoning with Think → Act → Observe → Reflect cycle",
    "✅ 5+ Code Analysis Tools — Fetch, structure analysis, security, performance, quality scoring",
    "✅ LLM Decision Making — Gemini 2.0 Flash (primary) with OpenAI as alternative provider",
    "✅ REST API + Web UI — FastAPI backend with auto-docs, Streamlit frontend with live feedback",
    "✅ Database + Testing — SQLite persistence, 257 tests, 85%+ coverage, 85/100 audit score",
]
add_bullet_slide(slide, "Features Implemented", features, "SPEAKER NOTES (1.5 min): Let me walk through what we actually built... The ReACT loop is the core innovation. We have 5 specialized tools. The system supports multiple LLM providers. Full API documentation is auto-generated by FastAPI. And we validated everything with 257 tests.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "Technology Stack", 36, WHITE, True)
line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), BLUE_ACCENT)

stacks = [
    ("Backend", [
        "Python 3.10+ — Core language",
        "FastAPI — REST API framework",
        "SQLAlchemy — ORM with SQLite",
        "Uvicorn — ASGI server",
    ], GREEN),
    ("Frontend", [
        "Streamlit — Web UI framework",
        "Requests — API client library",
        "Markdown — Rich text rendering",
    ], BLUE_ACCENT),
    ("AI / LLM", [
        "Gemini 2.0 Flash — Primary LLM",
        "OpenAI GPT-4o-mini — Alternative",
        "Google Generative AI SDK",
        "Custom prompt engineering",
    ], GOLD),
    ("Testing & Quality", [
        "Pytest — Test framework",
        "pytest-asyncio — Async tests",
        "pytest-cov — Coverage reports",
        "Ruff — Linter & formatter",
        "MyPy — Type checking",
    ], RGBColor(0xE7, 0x4C, 0x3C)),
]
for i, (title, items, color) in enumerate(stacks):
    left = Inches(0.6 + i * 3.15)
    add_rect(slide, left, Inches(1.6), Inches(2.9), Inches(5.2), RGBColor(0x14, 0x2B, 0x44), RGBColor(0x2C, 0x4A, 0x6F))
    add_textbox(slide, left + Inches(0.15), Inches(1.7), Inches(2.6), Inches(0.4), title, 22, color, True)
    for j, item in enumerate(items):
        add_textbox(slide, left + Inches(0.15), Inches(2.3 + j * 0.5), Inches(2.6), Inches(0.5), f"▸ {item}", 14, LIGHT_GRAY)

notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (1 min): Our tech stack is modern Python ecosystem... FastAPI for backend, Streamlit for frontend, Gemini for AI reasoning. We also have comprehensive testing infrastructure with Ruff linting and MyPy type checking."

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: PROJECT STATISTICS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "Project Statistics", 36, WHITE, True)
line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), BLUE_ACCENT)

stats = [
    ("6,000+", "Lines of Code", GREEN),
    ("27", "Source Files", BLUE_ACCENT),
    ("257", "Tests Passing", GOLD),
    ("85+%", "Code Coverage", RGBColor(0xE7, 0x4C, 0x3C)),
    ("85/100", "Audit Score", RGBColor(0x8E, 0x44, 0xAD)),
    ("2 Weeks", "Development", GREEN),
]
for i, (value, label, color) in enumerate(stats):
    left = Inches(0.6 + (i % 3) * 4.1)
    top = Inches(1.6 + (i // 3) * 2.8)
    add_rect(slide, left, top, Inches(3.7), Inches(2.3), RGBColor(0x14, 0x2B, 0x44), RGBColor(0x2C, 0x4A, 0x6F))
    add_textbox(slide, left + Inches(0.2), top + Inches(0.3), Inches(3.3), Inches(1.0), value, 44, color, True, PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), top + Inches(1.4), Inches(3.3), Inches(0.5), label, 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (1 min): The numbers speak for themselves... 6,000+ lines of production Python code, 257 tests all passing, 85%+ code coverage, and an 85/100 audit score from our internal code review process."

# ═══════════════════════════════════════════════════════════════
# SLIDE 12: HOW IT WORKS (DEMO FLOW)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "How It Works — Demo Flow", 36, WHITE, True)
line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), BLUE_ACCENT)

steps = [
    ("1", "Submit URL", "Paste GitHub\nrepo URL in UI"),
    ("2", "Agent Starts", "ReACT loop\ninitializes context"),
    ("3", "Fetch Code", "Download &\nanalyze structure"),
    ("4", "Run Tools", "Security, perf,\nquality checks"),
    ("5", "Generate", "LLM creates\nreview report"),
    ("6", "Display", "Results shown\nin Streamlit UI"),
]
for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.6 + i * 2.1)
    add_rect(slide, left, Inches(1.6), Inches(1.85), Inches(0.6), BLUE_ACCENT)
    add_textbox(slide, left + Inches(0.05), Inches(1.65), Inches(1.75), Inches(0.5), f"Step {num}: {title}", 16, WHITE, True, PP_ALIGN.CENTER)
    add_rect(slide, left, Inches(2.3), Inches(1.85), Inches(1.5), RGBColor(0x14, 0x2B, 0x44), RGBColor(0x2C, 0x4A, 0x6F))
    add_textbox(slide, left + Inches(0.1), Inches(2.4), Inches(1.65), Inches(1.3), desc, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_arrow_right(slide, left + Inches(1.9), Inches(1.75))

add_textbox(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.4), "Total time: ~30-60 seconds per review • Results include: quality score, issues found, security vulnerabilities, performance bottlenecks", 16, GOLD, False, PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (1.5 min): Let me walk through the demo flow... User submits a GitHub URL. Agent starts the ReACT loop. It fetches the code, runs analysis tools, generates a report using LLM, and displays results. The whole process takes 30-60 seconds."

# ═══════════════════════════════════════════════════════════════
# SLIDE 13: RESULTS / SCREENSHOTS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "Results & Demo", 36, WHITE, True)
line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), BLUE_ACCENT)

add_image_placeholder(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.5), "[ Streamlit UI Screenshot ]")
add_image_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5), "[ API Docs Screenshot ]")
add_image_placeholder(slide, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.5), "[ Review Results Screenshot ]")
add_image_placeholder(slide, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.5), "[ Test Coverage Screenshot ]")

notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (1.5 min): Here are screenshots of the working system... The Streamlit UI is clean and professional. FastAPI automatically generates API docs. Review results show quality scores with detailed issue breakdown. Our test coverage is 85%+ across all modules."

# ═══════════════════════════════════════════════════════════════
# SLIDE 14: LESSONS LEARNED
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Lessons Learned", [
    "ReACT framework is powerful but prompt engineering is critical for quality results",
    "Tool design matters — well-designed tools make the agent smarter and more reliable",
    "Testing saves time — 257 tests caught regressions early and gave confidence to refactor",
    "Clear code structure helps — modular design let 5 people work in parallel without conflicts",
    "Team coordination is essential — daily standups and shared docs kept everyone aligned",
], "SPEAKER NOTES (1.5 min): We learned a lot building this... The ReACT framework is powerful but the quality depends heavily on how you prompt the LLM. Well-designed tools make the agent smarter. Testing caught many regressions. Modular code let us work in parallel.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 15: FUTURE IMPROVEMENTS & CONCLUSION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "Future Improvements & Conclusion", 36, WHITE, True)
line = add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(0.04), BLUE_ACCENT)

add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.4), "Future Improvements", 24, GOLD, True)
futures = [
    "Multi-language support (Java, C++, Go)",
    "Custom rule definitions for teams",
    "Automated PR comments & fix suggestions",
    "CI/CD pipeline integration",
    "Real-time collaborative reviews",
]
for i, f in enumerate(futures):
    add_textbox(slide, Inches(0.6), Inches(2.0 + i * 0.45), Inches(5.5), Inches(0.4), f"🔹 {f}", 16, LIGHT_GRAY)

add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4), "Conclusion", 24, GREEN, True)
conclusions = [
    "Successfully built an agentic AI code review system",
    "ReACT loop architecture proven effective",
    "85/100 audit score validates code quality",
    "257 tests ensure reliability",
    "Team collaborated effectively across all roles",
]
for i, c in enumerate(conclusions):
    add_textbox(slide, Inches(7.0), Inches(2.0 + i * 0.45), Inches(5.5), Inches(0.4), f"✅ {c}", 16, LIGHT_GRAY)

notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (1.5 min): Looking ahead, we could add multi-language support, custom rules, and CI/CD integration. But even in its current state, the system works well — 85/100 audit score, 257 tests, and a successful team collaboration across all roles."

# ═══════════════════════════════════════════════════════════════
# SLIDE 16: Q&A
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
acc = add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE_ACCENT)
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), "Questions & Discussion", 36, WHITE, True, PP_ALIGN.CENTER)
line = add_rect(slide, Inches(4.5), Inches(1.2), Inches(4), Inches(0.04), GOLD)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0), "Thank you for your time!", 40, GOLD, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6), "We welcome your questions and feedback", 24, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5), "Contact: code-review-agent@project.com", 16, RGBColor(0x8A, 0xAA, 0xC8), False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5), "Repository: github.com/code-review-agent", 16, RGBColor(0x8A, 0xAA, 0xC8), False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.4), "Slides available at: presentations/Code-Review-Agent-Overview.pptx", 14, RGBColor(0x6A, 0x8A, 0xA8), False, PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "SPEAKER NOTES (1 min): Thank you everyone for listening! We're happy to answer any questions about the architecture, implementation decisions, or how we'd extend this further."

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = "G:/Code-Review-Agent/project-github/presentations/Code-Review-Agent-Overview.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
